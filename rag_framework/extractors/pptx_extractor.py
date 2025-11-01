"""Extracteur basé sur python-pptx (extraction de présentations PowerPoint)."""

from pathlib import Path
from typing import Any, ClassVar

from rag_framework.extractors.base import BaseExtractor, ExtractionResult
from rag_framework.utils.logger import get_logger

logger = get_logger(__name__)


class PptxExtractor(BaseExtractor):
    """Extracteur utilisant python-pptx pour l'extraction de présentations PowerPoint.

    python-pptx est la librairie de référence en 2025 pour l'extraction
    de présentations Microsoft PowerPoint (.pptx, .pptm).
    Plus rapide et précis que Docling pour les présentations simples.

    Avantages:
    - Extraction native et rapide de .pptx
    - Accès aux titres, sous-titres et contenu des diapositives
    - Extraction de tableaux et listes
    - Extraction des notes de présentation
    - Support des métadonnées (auteur, date, etc.)
    - Pas de dépendance lourde

    Limitations:
    - Pas de support direct des .ppt (ancien format)
    - Images et graphiques non extraits (seulement texte)
    - Animations et transitions ignorées
    - SmartArt et diagrammes complexes simplifiés

    Parameters
    ----------
    config : dict[str, Any]
        Configuration de l'extracteur.
        Clés supportées:
        - extract_notes : bool (défaut: True)
        - extract_tables : bool (défaut: True)
        - include_slide_numbers : bool (défaut: True)
        - min_text_length : int (défaut: 10)
        - extract_metadata : bool (défaut: True)

    Notes
    -----
    Pour les présentations avec images complexes ou mise en page sophistiquée,
    Docling ou VLM peuvent être utilisés en fallback.
    """

    SUPPORTED_EXTENSIONS: ClassVar[set[str]] = {".pptx", ".pptm"}

    def can_extract(self, file_path: Path) -> bool:
        """Vérifie si le fichier est une présentation PowerPoint moderne.

        Parameters
        ----------
        file_path : Path
            Chemin vers le fichier.

        Returns
        -------
        bool
            True si le fichier a l'extension .pptx ou .pptm.

        Examples
        --------
        >>> extractor = PptxExtractor(config={})
        >>> extractor.can_extract(Path("presentation.pptx"))
        True
        >>> extractor.can_extract(Path("presentation.ppt"))
        False
        """
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(self, file_path: Path) -> ExtractionResult:
        """Extrait le texte d'une présentation PowerPoint avec python-pptx.

        Parameters
        ----------
        file_path : Path
            Chemin vers le fichier PowerPoint.

        Returns
        -------
        ExtractionResult
            Résultat de l'extraction.

        Notes
        -----
        L'extraction traite chaque diapositive en incluant:
        - Titre et sous-titre
        - Contenu des zones de texte
        - Tableaux (si extract_tables=True)
        - Notes de présentation (si extract_notes=True)
        """
        try:
            # Import tardif pour éviter erreur si librairie non installée
            from pptx import Presentation

            # Options d'extraction
            extract_notes = self.config.get("extract_notes", True)
            extract_tables = self.config.get("extract_tables", True)
            include_slide_numbers = self.config.get("include_slide_numbers", True)

            # Ouverture de la présentation
            prs = Presentation(str(file_path))

            text_parts = []
            slide_count = 0
            table_count = 0

            # Extraction diapositive par diapositive
            for slide_idx, slide in enumerate(prs.slides):
                slide_count += 1
                slide_num = slide_idx + 1

                if include_slide_numbers:
                    text_parts.append(f"\n{'=' * 60}")
                    text_parts.append(f"DIAPOSITIVE {slide_num}")
                    text_parts.append(f"{'=' * 60}\n")

                # Extraction du texte de chaque forme (shape)
                for shape in slide.shapes:
                    # Titre ou zone de texte
                    if hasattr(shape, "text") and shape.text.strip():
                        # Détection du titre
                        if shape.is_placeholder:
                            placeholder = shape.placeholder_format
                            if placeholder.type == 1:  # PP_PLACEHOLDER.TITLE
                                text_parts.append(f"# {shape.text.strip()}")
                            elif (
                                placeholder.type == 3
                            ):  # PP_PLACEHOLDER.CENTER_TITLE
                                text_parts.append(f"# {shape.text.strip()}")
                            elif placeholder.type == 2:  # PP_PLACEHOLDER.SUBTITLE
                                text_parts.append(f"## {shape.text.strip()}")
                            else:
                                text_parts.append(shape.text.strip())
                        else:
                            text_parts.append(shape.text.strip())

                    # Tableau
                    elif extract_tables and shape.has_table:
                        table_count += 1
                        table_text = self._extract_table(
                            shape.table, slide_num, table_count
                        )
                        if table_text:
                            text_parts.append(table_text)

                # Notes de présentation
                if extract_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        text_parts.append(f"\n**Notes:**\n{notes_text}")

            # Concaténation
            full_text = "\n\n".join(text_parts)

            # Métadonnées
            metadata: dict[str, Any] = {
                "file_size": file_path.stat().st_size,
                "file_name": file_path.name,
                "extractor": "python-pptx",
                "slides_count": slide_count,
                "tables_count": table_count,
            }

            # Extraction des métadonnées de la présentation
            if self.config.get("extract_metadata", True):
                core_props = prs.core_properties
                if core_props.title:
                    metadata["title"] = core_props.title
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.subject:
                    metadata["subject"] = core_props.subject
                if core_props.keywords:
                    metadata["keywords"] = core_props.keywords
                if core_props.created:
                    metadata["creation_date"] = str(core_props.created)
                if core_props.modified:
                    metadata["modification_date"] = str(core_props.modified)

            # Vérification de la longueur minimale
            min_length = self.config.get("min_text_length", 10)
            if len(full_text.strip()) < min_length:
                return ExtractionResult(
                    text=full_text,
                    success=False,
                    extractor_name=self.name,
                    metadata=metadata,
                    error=f"Texte extrait trop court ({len(full_text)} < {min_length})",
                    confidence_score=0.1,
                )

            # Score de confiance
            # python-pptx est très fiable pour les .pptx natifs
            confidence = 0.95 if len(full_text) > 50 else 0.7

            logger.debug(
                f"python-pptx: Extrait {len(full_text)} caractères "
                f"({slide_count} diapositives, {table_count} tableaux) "
                f"(confidence={confidence:.2f})"
            )

            return ExtractionResult(
                text=full_text,
                success=True,
                extractor_name=self.name,
                metadata=metadata,
                confidence_score=confidence,
            )

        except ImportError:
            error_msg = (
                "python-pptx n'est pas installé. "
                "Installez avec: pip install python-pptx"
            )
            logger.error(error_msg)
            return ExtractionResult(
                text="",
                success=False,
                extractor_name=self.name,
                metadata={},
                error=error_msg,
                confidence_score=0.0,
            )

        except Exception as e:
            error_msg = f"Erreur python-pptx extraction: {e}"
            logger.warning(error_msg)
            return ExtractionResult(
                text="",
                success=False,
                extractor_name=self.name,
                metadata={},
                error=error_msg,
                confidence_score=0.0,
            )

    def _extract_table(self, table: Any, slide_num: int, table_num: int) -> str:
        """Extrait le contenu d'un tableau au format Markdown.

        Parameters
        ----------
        table : Table
            Objet tableau de python-pptx.
        slide_num : int
            Numéro de la diapositive.
        table_num : int
            Numéro du tableau dans la présentation.

        Returns
        -------
        str
            Tableau formaté en Markdown.
        """
        rows = []
        rows.append(f"\n### Tableau {table_num} (Diapositive {slide_num})\n")

        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            # Format Markdown
            rows.append("| " + " | ".join(cells) + " |")

            # Ligne de séparation après l'en-tête (première ligne)
            if row_idx == 0 and len(cells) > 0:
                separator = "|" + "|".join(["---"] * len(cells)) + "|"
                rows.append(separator)

        return "\n".join(rows)
