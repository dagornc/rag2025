"""Adapter python-pptx pour les fichiers PowerPoint (.pptx).

Auteur: RAG Framework Team
Version: 1.0.0
"""

from typing import Any, ClassVar

from rag_framework.preprocessing.adapters.base import LibraryAdapter, ParsingError


class PythonPptxAdapter(LibraryAdapter):
    """Adapter pour la librairie python-pptx.

    Extrait le texte des fichiers PowerPoint (.pptx).

    Attributes:
        REQUIRED_MODULES: Liste des modules requis.
    """

    REQUIRED_MODULES: ClassVar[list[str]] = ["pptx"]

    def parse(self, file_path: str) -> dict[str, Any]:
        """Parse un fichier PowerPoint avec python-pptx.

        Args:
            file_path: Chemin vers le fichier .pptx.

        Returns:
            Dictionnaire avec text, metadata, slides.

        Raises:
            ParsingError: Si le parsing échoue.
        """
        try:
            from pptx import Presentation

            # Ouvrir la présentation
            prs = Presentation(file_path)

            # Configuration
            extract_notes = self.library_config.get("extract_speaker_notes", True)

            # Extraire le texte de toutes les slides
            slides = []
            full_text = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_dict: dict[str, Any] = {
                    "slide_number": slide_num,
                    "text": "",
                    "shapes": [],
                }

                # Extraire le texte de toutes les formes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                        slide_dict["shapes"].append(
                            {"text": shape.text, "type": shape.shape_type.name}
                        )

                # Extraire les notes du présentateur
                if extract_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_dict["notes"] = notes_text
                        slide_text.append(f"[Notes: {notes_text}]")

                slide_dict["text"] = "\n".join(slide_text)
                slides.append(slide_dict)
                full_text.extend(slide_text)

            # Métadonnées
            core_props = prs.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "slide_count": len(slides),
            }

            return {
                "text": "\n\n".join(full_text),
                "metadata": metadata,
                "slides": slides,
            }

        except Exception as e:
            raise ParsingError(f"Échec python-pptx parsing : {e}") from e
